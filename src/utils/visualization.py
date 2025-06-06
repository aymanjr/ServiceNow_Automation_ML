import argparse
import json
from pathlib import Path
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def map_service_name(name):
    mapping = {
        "AUTO": "Terminal Automation",
        "AD": "Asset Digitalisation",
        "IW": "Industrial Wireless",
        "CF": "Confluence",
        "MVM": "Mavim",
        "HV": "Horizon View"
    }
    return mapping.get(name.upper(), name)

def add_summary_slide(prs, data, overall):
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    table_left = Inches(5.65)
    table_top = Inches(3.05)
    table_width = Inches(7.47)
    table_height = Inches(4.04)
    table = slide.shapes.add_table(11, 4, table_left, table_top, table_width, table_height).table

    headers = ["Services", "Incidents", "Requests", "Total Tickets"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        table.cell(0, i).text_frame.paragraphs[0].font.bold = True

    order = ["SIP", "FLOW", "AUTO", "AD", "IW", "CF", "MVM", "HV", "IFS"]
    total_inc = total_ritm = total_all = 0

    for i, s in enumerate(order):
        stats = data.get(s.upper(), {})
        inc = int(stats.get("INC_count", 0) or 0)
        ritm = int(stats.get("RITM_count", 0) or 0)
        total = inc + ritm
        total_inc += inc
        total_ritm += ritm
        total_all += total
        table.cell(i + 1, 0).text = map_service_name(s)
        table.cell(i + 1, 1).text = str(inc)
        table.cell(i + 1, 2).text = str(ritm)
        table.cell(i + 1, 3).text = str(total)

    table.cell(10, 0).text = "Total"
    table.cell(10, 1).text = str(total_inc)
    table.cell(10, 2).text = str(total_ritm)
    table.cell(10, 3).text = str(total_all)
    for i in range(4):
        table.cell(10, i).text_frame.paragraphs[0].font.bold = True

    if Path("data/charts/urgency_heatmap_INC.png").exists():
        slide.shapes.add_picture("data/charts/urgency_heatmap_INC.png", Inches(10.51), Inches(0.43), width=Inches(2.77), height=Inches(2.15))

    if Path("data/charts/urgency_heatmap_RITM.png").exists():
        slide.shapes.add_picture("data/charts/urgency_heatmap_RITM.png", Inches(5.44), Inches(0.48), width=Inches(2.64), height=Inches(2.05))

    if Path("data/charts/donut_total.png").exists():
        slide.shapes.add_picture("data/charts/donut_total.png", Inches(8.03), Inches(0.43), width=Inches(2.51), height=Inches(2.28))

    keynotes_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.76), Inches(2.68), Inches(4.03), Inches(3.75))
    keynotes_box.fill.solid()
    keynotes_box.fill.fore_color.rgb = RGBColor(255, 165, 0)
    keynotes_box.text = "Weekly Keynotes (to be filled manually)"
    keynotes_box.text_frame.paragraphs[0].font.size = Pt(14)
    keynotes_box.text_frame.paragraphs[0].font.bold = True

def add_insights_slide(prs, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.text = "Key Weekly Insights"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True

    comparison = summary.get("comparison", {})
    insights_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(6), Inches(3.5))
    insights_tf = insights_box.text_frame
    insights_tf.word_wrap = True
    for key, line in comparison.items():
        p = insights_tf.add_paragraph()
        p.text = f"- {line}"
        p.font.size = Pt(14)
    if insights_tf.paragraphs:
        insights_tf.paragraphs[0].font.bold = True

    if Path("data/charts/volume_by_service.png").exists():
        slide.shapes.add_picture("data/charts/volume_by_service.png", Inches(7), Inches(1.1), height=Inches(3.0))

    if Path("data/charts/monthly_progress.png").exists():
        slide.shapes.add_picture("data/charts/monthly_progress.png", Inches(0.5), Inches(4.7), height=Inches(2.5))

def generate_ppt(json_path: str, output_dir: str):
    with open(json_path, 'r') as f:
        summary = json.load(f)

    data = summary.get("services", {})
    overall = summary.get("overall", {})
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    add_summary_slide(prs, data, overall)
    add_insights_slide(prs, summary)

    ppt_path = output_path / "Service_Report.pptx"
    prs.save(ppt_path)
    print(f"[OK] Presentation saved to {ppt_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate service report PowerPoint")
    parser.add_argument("--input", required=True, help="Path to JSON summary file")
    parser.add_argument("--output", required=True, help="Directory to save PPT")
    args = parser.parse_args()

    generate_ppt(args.input, args.output)
